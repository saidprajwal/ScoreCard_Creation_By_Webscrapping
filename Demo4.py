import os
from pptx import Presentation
import PyPDF2
import openpyxl

break_point = 1
folder="C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel"

wb = openpyxl.load_workbook("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")
sh2 = wb['Sheet2']

# store scores from excel in list
scores = []
for i in range(2, sh2.max_row + 1):
    if (sh2.cell(row=i, column=1).value == None):
        break_point = i
        break
    scores.append(sh2.cell(row=i, column=1).value)

# print(scores)
cols=2
def extract_pdf(folder,f):
    path=folder+'\\'+f
    reader = PyPDF2.PdfReader(path)
    pages = len(reader.pages)
    str = ""
    c = sh2.cell(row=1, column=cols)
    c.value=f
    wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")
    # check if score is present in pdf and store value
    total_score = 0
    rows = 2

    for item in scores:
        flag = 0
        c = sh2.cell(row=rows, column=cols)
        for i in range(0, pages):
            str = reader.pages[i].extract_text()
            if item.lower() in str.lower():
                flag=1
                c.value = 1
                total_score += 1
                wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")
                break
        rows = rows + 1
        if flag==0:
              c.value=0
              wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")


    # store total score obtained in break_point row
    #c = sh2.cell(row=break_point, column=cols)
    #c.value = "{}/{}".format(total_score, len(scores))
    c = sh2.cell(row=break_point, column=cols)
    c.value = (total_score / len(scores)) * 100
    wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")


def extract_ppt(folder,f):
    path = folder + '\\' + f
    prs = Presentation(path)
    c = sh2.cell(row=1, column=cols)
    c.value = f
    wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")
    str = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    str = str + run.text


    # store scores from excel in list

    total_score = 0
    rows = 2
    for item in scores:
        c = sh2.cell(row=rows, column=cols)
        if item.lower() in str.lower():
            c.value = 1;
            total_score += 1;
            wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")
        else:
            c.value = 0
            wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")
        rows = rows + 1


    c = sh2.cell(row=break_point, column=cols)
    c.value = (total_score / len(scores)) * 100
    wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")


for f in os.listdir(folder):
    if f.endswith(".pdf"):
        extract_pdf(folder, f)
        cols = cols + 1
    elif f.endswith(".pptx"):
        extract_ppt(folder,f)
        cols = cols + 1