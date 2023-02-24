from pptx import Presentation
import openpyxl

wb=openpyxl.load_workbook("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")
sh2=wb['Sheet2']

prs = Presentation("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Obj.pptx")
str=""

for slide in prs.slides:
     for shape in slide.shapes:
         if not shape.has_text_frame:
             continue
         for paragraph in shape.text_frame.paragraphs:
             for run in paragraph.runs:
                 str=str+run.text
print(str)

break_point=1
#store scores from excel in list
scores=[]
for i in range (2,sh2.max_row+1):
    if(sh2.cell(row=i,column=1).value==None):
         break_point=i
         break
    scores.append(sh2.cell(row=i,column=1).value)

total_score=0
rows=2
for item in scores:
    c = sh2.cell(row=rows, column=5)
    if item.lower() in str.lower():
        c.value = 1;
        total_score += 1;
        wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")
    rows = rows + 1

c=sh2.cell(row=break_point,column=5)
c.value="{}/{}".format(total_score,len(scores))
wb.save("C:\\Users\\PRAJWAL\\PycharmProjects\\WebScapping\\excel\\Audit.xlsx")

